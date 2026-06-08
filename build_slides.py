"""Fill the project presentation template with real content + figures.

Edits Project_Presentation_Template.pptx in place (6 slides matching the
required structure: Title, Intro, Dataset, Model, Results, Conclusions).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DECK = "Project_Presentation_Template.pptx"
prs = Presentation(DECK)
slides = list(prs.slides)


def set_title(slide, text):
    slide.shapes.title.text = text
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True


def get_body(slide):
    """The 'Content Placeholder 2' on a Title-and-Content slide."""
    for sh in slide.shapes:
        if sh.is_placeholder and sh.name.startswith("Content"):
            return sh
    raise KeyError("no content placeholder")


def set_bullets(slide, items):
    """items: list of (text, level, bold, size_pt)."""
    tf = get_body(slide).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (text, level, bold, size) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def replace_picture(slide, img_path, fit="height"):
    """Swap the existing 'Picture 3' for img_path, centered in its box."""
    from PIL import Image
    pic = next(sh for sh in slide.shapes if sh.name.startswith("Picture"))
    box_l, box_t, box_w, box_h = pic.left, pic.top, pic.width, pic.height
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    if fit == "height":
        h = box_h
        w = int(h * ar)
        if w > box_w:
            w = box_w
            h = int(w / ar)
    else:
        w = box_w
        h = int(w / ar)
        if h > box_h:
            h = box_h
            w = int(h * ar)
    l = box_l + (box_w - w) // 2
    t = box_t + (box_h - h) // 2
    pic._element.getparent().remove(pic._element)
    slide.shapes.add_picture(img_path, l, t, w, h)


def remove_shape(slide, name):
    for sh in list(slide.shapes):
        if sh.name == name:
            sh._element.getparent().remove(sh._element)


def set_caption(slide, name, text):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            sh.text_frame.paragraphs[0].text = text


# --- Slide 0: Title -------------------------------------------------------
slides[0].shapes.title.text = "Visual Identification of Electric Guitars"
slides[0].placeholders[1].text = (
    "Classifying guitar photos by make & model  —  Stratocaster vs. Telecaster\n"
    "Jeremy Mitaux"
)
set_notes(slides[0],
    "Hi, I'm Jeremy. My project trains a deep-learning model to recognize what "
    "kind of guitar is in a photo — specifically a Fender Stratocaster vs. a "
    "Telecaster — straight from a marketplace listing image.")

# --- Slide 1: Introduction ------------------------------------------------
set_title(slides[1], "Introduction")
set_bullets(slides[1], [
    ("What is my project", 0, True, 20),
    ("Predict a guitar's make/model from one photo", 1, False, 18),
    ("A CNN image classifier: Fender Stratocaster vs. Telecaster", 1, False, 18),
    ("Why I'm interested", 0, True, 20),
    ("I want an AI “flip-finding” assistant — spot mislabeled or "
     "underpriced guitar listings faster than reading every description", 1, False, 18),
    ("Combines a personal interest in guitars with deep learning", 1, False, 18),
])
set_notes(slides[1],
    "The long-term goal is a flip-finding assistant: point it at a listing photo "
    "and have it tell you what the guitar actually is, so you can catch listings "
    "that just say 'Fender Strat' or are mislabeled. The core learnable sub-problem "
    "is identifying make/model from an image — that's what I built. It's a real "
    "computer-vision task because different models have distinct silhouettes, "
    "headstocks, and pickup layouts.")

# --- Slide 2: Dataset -----------------------------------------------------
set_title(slides[2], "Dataset")
set_bullets(slides[2], [
    ("What data", 0, True, 20),
    ("6,195 photos from 1,333 real Reverb.com listings", 1, False, 18),
    ("862 Stratocaster / 471 Telecaster — built from scratch, no existing dataset", 1, False, 18),
    ("How I built it", 0, True, 20),
    ("Scraped Reverb via the Apify API; labeled by model from listing metadata", 1, False, 18),
    ("Filtered out parts, wrong brands, cheap accessories (dropped 278 listings)", 1, False, 18),
    ("Split grouped by listing — no guitar appears in both train and test", 1, False, 18),
])
replace_picture(slides[2], "data/sample_images/stratocaster_10366011_0.jpg", fit="height")
set_caption(slides[2], "TextBox 4", "A Stratocaster listing photo from the dataset")
set_notes(slides[2],
    "There's no off-the-shelf dataset, so I built one from live Reverb listings. "
    "I scraped per-model, then labeled each listing by its structured make/model "
    "metadata rather than guessing from keywords, and filtered out parts, "
    "wrong-brand sellers, and cheap accessories. One important detail: each listing "
    "has up to 5 near-duplicate photos, so I split the data BY LISTING — every photo "
    "of one guitar stays in a single split. A naive per-image split would leak "
    "duplicates into both train and test and fake a high accuracy.")

# --- Slide 3: Model -------------------------------------------------------
set_title(slides[3], "Your Model")
set_bullets(slides[3], [
    ("What model", 0, True, 20),
    ("EfficientNet-B0 CNN, pretrained on ImageNet", 1, False, 18),
    ("Fine-tuned with a 2-way head (AdamW, cosine LR, label smoothing)", 1, False, 18),
    ("Why I picked it", 0, True, 20),
    ("Strong accuracy at a small size — trains fast on a laptop GPU (Apple MPS)", 1, False, 18),
    ("Transfer learning works well with a modest ~1,300-listing dataset", 1, False, 18),
    ("Body shape / headstock / pickups are real signals a CNN can learn", 1, False, 18),
])
set_notes(slides[3],
    "I used EfficientNet-B0, a convolutional neural network pretrained on ImageNet, "
    "and fine-tuned it for two classes. I picked it because it hits a great accuracy-"
    "to-size ratio — it trains quickly on my Mac's GPU — and transfer learning lets a "
    "pretrained backbone do well even with only ~1,300 listings. I also handled class "
    "imbalance with a weighted sampler and used data augmentation to reduce overfitting.")

# --- Slide 4: Results -----------------------------------------------------
set_title(slides[4], "Your Results")
set_bullets(slides[4], [
    ("Headline metric", 0, True, 20),
    ("95.2% test accuracy on 933 held-out images (200 listings)", 1, False, 18),
    ("Macro-F1 0.95  —  vs. a 64.2% majority-class baseline", 1, False, 18),
    ("Per-class F1: Stratocaster 0.96, Telecaster 0.93", 1, False, 18),
    ("How did it work", 0, True, 20),
    ("Worked well — but Strat vs. Tele is a visually easy 2-class problem", 1, False, 18),
    ("Read it as a strong proof of concept, not a final number", 1, False, 18),
])
# The pre-labeled confusion matrix is self-contained; drop the template's
# redundant axis-label textboxes and keep one caption.
remove_shape(slides[4], "TextBox 4")  # "Predicted Values"
remove_shape(slides[4], "TextBox 5")  # "True Values"
replace_picture(slides[4], "results/confusion_matrix.png", fit="height")
set_caption(slides[4], "TextBox 6", "Confusion matrix — held-out test set")
set_notes(slides[4],
    "The model got 95.2% accuracy on a held-out test set of 933 images. I report a "
    "few things for honesty: macro-F1 of 0.95, and the majority-class baseline of "
    "64.2% — that's what you'd get always guessing 'Stratocaster' — so you can see the "
    "real lift. The confusion matrix shows errors are balanced, not one class "
    "collapsing into the other. The big caveat: Strat vs. Tele is a visually easy "
    "distinction, so this is a strong result on an EASY problem — I'd expect it to "
    "drop as look-alike models get added.")

# --- Slide 5: Conclusions -------------------------------------------------
set_title(slides[5], "Conclusions")
set_bullets(slides[5], [
    ("What worked", 0, True, 20),
    ("Full pipeline works end-to-end: scrape → label → train → evaluate, reproducible", 1, False, 18),
    ("Group-aware split gives an honest accuracy (no photo leakage)", 1, False, 18),
    ("What I changed from the proposal", 0, True, 20),
    ("Dropped year (only 44% of listings had one, and it was messy free-text)", 1, False, 18),
    ("Pivoted from country-of-origin (Strats look identical) to make/model — a real visual signal", 1, False, 18),
    ("What I'd do next", 0, True, 20),
    ("Scale to more models (Les Paul, SG, ES-335) — pipeline already supports N classes", 1, False, 18),
])
set_notes(slides[5],
    "Overall it worked like I expected — the pipeline is fully reproducible and the "
    "group-aware split means the accuracy is honest. Two things I'd flag: I changed "
    "scope from the proposal. I dropped predicting the year, because only 44% of "
    "listings had one and it was messy free text. And I pivoted away from guessing "
    "country of origin — American vs Mexican Strats look identical, so that model was "
    "really learning backgrounds, not guitars. Make/model is a genuine visual signal. "
    "If I kept going, I'd add more models like the Les Paul, SG, and ES-335 — the code "
    "already scales to N classes; I stopped at two when my scraping credits ran out. "
    "That harder problem is the more interesting next step.")

prs.save(DECK)
print(f"Saved {DECK} with {len(prs.slides)} slides.")
